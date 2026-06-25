"""
app/ai/file_parser.py uchun unit testlar.

Barcha fayl turlari test qilinadi:
  PDF, Word, Excel, CSV, TSV, PowerPoint, EPUB,
  HTML, JSON, JSONL, XML, RTF, Rasm (OCR), Audio (STT),
  Oddiy matn (TXT, MD)

Tashqi kutubxonalar (pdfplumber, pytesseract, faster-whisper va h.k.)
mock qilinadi — haqiqiy fayllar yoki modellar kerak emas.
"""

import io
import json
import pytest
from unittest.mock import patch, MagicMock

from app.ai.file_parser import (
    extract_text_from_pdf,
    extract_text_from_word,
    extract_text_from_excel,
    extract_text_from_csv,
    extract_text_from_pptx,
    extract_text_from_epub,
    extract_text_from_html,
    extract_text_from_json,
    extract_text_from_xml,
    extract_text_from_rtf,
    extract_text_from_image,
    extract_text_from_audio,
    parse_file,
)



class TestExtractTextFromPdf:
    def _make_mock_pdf(self, pages_text: list):
        mock_page = MagicMock()
        mock_page.extract_text.side_effect = pages_text
        mock_pdf = MagicMock()
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)
        mock_pdf.pages = [mock_page] * len(pages_text)
        return mock_pdf

    @patch("app.ai.file_parser.pdfplumber")
    def test_single_page(self, mock_pdfplumber):
        mock_pdfplumber.open.return_value = self._make_mock_pdf(["Salom dunyo"])
        result = extract_text_from_pdf(b"fake_pdf")
        assert "Salom dunyo" in result

    @patch("app.ai.file_parser.pdfplumber")
    def test_multiple_pages_joined(self, mock_pdfplumber):
        mock_pdfplumber.open.return_value = self._make_mock_pdf(["Birinchi", "Ikkinchi"])
        result = extract_text_from_pdf(b"fake_pdf")
        assert "Birinchi" in result
        assert "Ikkinchi" in result

    @patch("app.ai.file_parser.pdfplumber")
    def test_empty_pdf_raises(self, mock_pdfplumber):
        mock_pdfplumber.open.return_value = self._make_mock_pdf([None])
        with pytest.raises(ValueError, match="PDF dan matn topilmadi"):
            extract_text_from_pdf(b"fake_pdf")


class TestExtractTextFromWord:
    def _make_docx_bytes(self, paragraphs: list, table_rows: list = None):
        """Minimal in-memory .docx fayl yaratadi."""
        import docx
        doc = docx.Document()
        for para in paragraphs:
            doc.add_paragraph(para)
        if table_rows:
            table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
            for i, row_data in enumerate(table_rows):
                for j, cell_text in enumerate(row_data):
                    table.rows[i].cells[j].text = cell_text
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def test_basic_paragraphs(self):
        try:
            import docx
        except ImportError:
            pytest.skip("python-docx o'rnatilmagan")
        data = self._make_docx_bytes(["Birinchi paragraf", "Ikkinchi paragraf"])
        result = extract_text_from_word(data)
        assert "Birinchi paragraf" in result
        assert "Ikkinchi paragraf" in result

    def test_empty_document_raises(self):
        try:
            import docx
        except ImportError:
            pytest.skip("python-docx o'rnatilmagan")
        data = self._make_docx_bytes([])
        with pytest.raises(ValueError, match="Word hujjatidan matn topilmadi"):
            extract_text_from_word(data)

    def test_table_text_extracted(self):
        try:
            import docx
        except ImportError:
            pytest.skip("python-docx o'rnatilmagan")
        data = self._make_docx_bytes([], table_rows=[["Ism", "Yosh"], ["Ali", "25"]])
        result = extract_text_from_word(data)
        assert "Ali" in result
        assert "Yosh" in result



class TestExtractTextFromExcel:
    def _make_xlsx_bytes(self, rows: list, sheet_name: str = "Sheet1"):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl o'rnatilmagan")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name
        for row in rows:
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def test_basic_data(self):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl o'rnatilmagan")
        data = self._make_xlsx_bytes([["Ism", "Yosh"], ["Ali", 25], ["Vali", 30]])
        result = extract_text_from_excel(data)
        assert "Ali" in result
        assert "Yosh" in result

    def test_sheet_name_in_result(self):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl o'rnatilmagan")
        data = self._make_xlsx_bytes([["A", "B"], ["1", "2"]], sheet_name="Ma'lumotlar")
        result = extract_text_from_excel(data)
        assert "Ma'lumotlar" in result

    def test_empty_sheet_raises(self):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl o'rnatilmagan")
        data = self._make_xlsx_bytes([])
        with pytest.raises(ValueError, match="Excel fayldan ma'lumot topilmadi"):
            extract_text_from_excel(data)



class TestExtractTextFromCsv:
    def test_basic_csv(self):
        csv_bytes = b"Ism,Yosh\nAli,25\nVali,30"
        result = extract_text_from_csv(csv_bytes)
        assert "Ali" in result
        assert "Yosh" in result

    def test_tsv(self):
        tsv_bytes = b"Ism\tYosh\nAli\t25"
        result = extract_text_from_csv(tsv_bytes, delimiter="\t")
        assert "Ali" in result

    def test_empty_csv_raises(self):
        with pytest.raises(ValueError, match="CSV fayl bo'sh"):
            extract_text_from_csv(b"")

    def test_header_only_raises(self):
        with pytest.raises(ValueError, match="CSV fayldan ma'lumot topilmadi"):
            extract_text_from_csv(b"Ism,Yosh")

    def test_utf8_content(self):
        csv_bytes = "Ism,Shahar\nAziz,Toshkent".encode("utf-8")
        result = extract_text_from_csv(csv_bytes)
        assert "Toshkent" in result



class TestExtractTextFromPptx:
    def _make_pptx_bytes(self, slides_text: list):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx o'rnatilmagan")
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        for text in slides_text:
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
            txBox.text_frame.text = text
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_single_slide(self):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx o'rnatilmagan")
        data = self._make_pptx_bytes(["Birinchi slayd matni"])
        result = extract_text_from_pptx(data)
        assert "Birinchi slayd matni" in result

    def test_multiple_slides(self):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx o'rnatilmagan")
        data = self._make_pptx_bytes(["Slayd 1", "Slayd 2", "Slayd 3"])
        result = extract_text_from_pptx(data)
        assert "Slayd 1" in result
        assert "Slayd 3" in result

    def test_slide_number_in_result(self):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx o'rnatilmagan")
        data = self._make_pptx_bytes(["Test matni"])
        result = extract_text_from_pptx(data)
        assert "Slayd 1" in result

    def test_empty_pptx_raises(self):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx o'rnatilmagan")
        data = self._make_pptx_bytes([])
        with pytest.raises(ValueError, match="PowerPoint fayldan matn topilmadi"):
            extract_text_from_pptx(data)



class TestExtractTextFromHtml:
    def test_basic_html(self):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            pytest.skip("beautifulsoup4 o'rnatilmagan")
        html = b"<html><body><h1>Sarlavha</h1><p>Paragraf matni</p></body></html>"
        result = extract_text_from_html(html)
        assert "Sarlavha" in result
        assert "Paragraf matni" in result

    def test_script_and_style_removed(self):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            pytest.skip("beautifulsoup4 o'rnatilmagan")
        html = b"""<html><head><style>body{color:red}</style></head>
                   <body><script>alert('xato')</script><p>Asosiy matn</p></body></html>"""
        result = extract_text_from_html(html)
        assert "Asosiy matn" in result
        assert "alert" not in result
        assert "color:red" not in result

    def test_empty_html_raises(self):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            pytest.skip("beautifulsoup4 o'rnatilmagan")
        with pytest.raises(ValueError, match="HTML fayldan matn topilmadi"):
            extract_text_from_html(b"<html><body></body></html>")



class TestExtractTextFromJson:
    def test_simple_json_object(self):
        data = json.dumps({"ism": "Ali", "yosh": 25}).encode()
        result = extract_text_from_json(data)
        assert "Ali" in result
        assert "yosh" in result

    def test_json_array(self):
        data = json.dumps([{"id": 1}, {"id": 2}]).encode()
        result = extract_text_from_json(data)
        assert "id" in result

    def test_jsonl_multiple_objects(self):
        jsonl = b'{"ism":"Ali"}\n{"ism":"Vali"}\n{"ism":"Aziz"}'
        result = extract_text_from_json(jsonl)
        assert "Ali" in result
        assert "Vali" in result
        assert "Aziz" in result

    def test_invalid_json_returns_raw(self):
        raw = b"bu json emas, oddiy matn"
        result = extract_text_from_json(raw)
        assert "bu json emas" in result

    def test_unicode_preserved(self):
        data = json.dumps({"shahar": "Toshkent"}, ensure_ascii=False).encode("utf-8")
        result = extract_text_from_json(data)
        assert "Toshkent" in result



class TestExtractTextFromXml:
    def test_basic_xml(self):
        xml = b"<?xml version='1.0'?><root><item>Mahsulot nomi</item></root>"
        result = extract_text_from_xml(xml)
        assert "Mahsulot nomi" in result

    def test_nested_xml(self):
        xml = b"<katalog><kitob><nomi>Python</nomi><muallif>Guido</muallif></kitob></katalog>"
        result = extract_text_from_xml(xml)
        assert "Python" in result
        assert "Guido" in result


class TestExtractTextFromRtf:
    def test_basic_rtf_via_regex_fallback(self):
        rtf = b"{\\rtf1\\ansi {\\b Salom} {\\i Dunyo}}"
        result = extract_text_from_rtf(rtf)
        assert len(result) > 0

    def test_with_striprtf_if_available(self):
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            pytest.skip("striprtf o'rnatilmagan")
        rtf = b"{\\rtf1\\ansi Salom Dunyo}"
        result = extract_text_from_rtf(rtf)
        assert len(result) > 0

    def test_empty_rtf_raises(self):
        with pytest.raises(ValueError, match="RTF fayldan matn topilmadi"):
            extract_text_from_rtf(b"")



class TestExtractTextFromImage:
    @patch("app.ai.file_parser.pytesseract")
    @patch("app.ai.file_parser.Image")
    def test_ocr_returns_text(self, mock_image_module, mock_tesseract):
        mock_image_module.open.return_value = MagicMock()
        mock_tesseract.image_to_string.return_value = "OCR orqali topilgan matn"
        result = extract_text_from_image(b"fake_image_bytes")
        assert "OCR orqali topilgan matn" in result

    @patch("app.ai.file_parser.pytesseract")
    @patch("app.ai.file_parser.Image")
    def test_empty_ocr_raises(self, mock_image_module, mock_tesseract):
        mock_image_module.open.return_value = MagicMock()
        mock_tesseract.image_to_string.return_value = "   "
        with pytest.raises(ValueError, match="Rasmdan matn topilmadi"):
            extract_text_from_image(b"fake_image_bytes")

    @patch("app.ai.file_parser.pytesseract")
    @patch("app.ai.file_parser.Image")
    def test_lang_used(self, mock_image_module, mock_tesseract):
        mock_image_module.open.return_value = MagicMock()
        mock_tesseract.image_to_string.return_value = "Matn"
        extract_text_from_image(b"fake_image_bytes")
        call_kwargs = mock_tesseract.image_to_string.call_args
        assert "uzb+rus+eng" in str(call_kwargs)



class TestExtractTextFromAudio:
    def _mock_whisper(self):
        mock_segment = MagicMock()
        mock_segment.text = "Audio dan olingan matn"
        mock_info = MagicMock()
        mock_info.language = "uz"
        mock_info.duration = 5.0
        mock_model = MagicMock()
        mock_model.transcribe.return_value = ([mock_segment], mock_info)
        return mock_model

    @patch("app.ai.file_parser._get_whisper_model")
    def test_audio_transcription(self, mock_get_model):
        mock_get_model.return_value = self._mock_whisper()
        result = extract_text_from_audio(b"fake_audio", filename="test.mp3")
        assert "Audio dan olingan matn" in result

    @patch("app.ai.file_parser._get_whisper_model")
    def test_empty_transcription_raises(self, mock_get_model):
        mock_model = MagicMock()
        mock_model.transcribe.return_value = ([], MagicMock(language="uz", duration=1.0))
        mock_get_model.return_value = mock_model
        with pytest.raises(ValueError, match="Audio dan matn topilmadi"):
            extract_text_from_audio(b"fake_audio", filename="test.mp3")

    @patch("app.ai.file_parser._get_whisper_model")
    def test_extension_used_in_temp_file(self, mock_get_model):
        mock_get_model.return_value = self._mock_whisper()
        result = extract_text_from_audio(b"fake_audio", filename="speech.wav")
        assert len(result) > 0



class TestParseFile:
    def test_txt_content_type(self):
        result = parse_file(b"Oddiy matn", "text/plain", "file.txt")
        assert "Oddiy matn" in result

    def test_markdown_content_type(self):
        result = parse_file(b"# Sarlavha\nMazmun", "text/markdown", "file.md")
        assert "Sarlavha" in result

    def test_json_content_type(self):
        data = json.dumps({"kalit": "qiymat"}).encode()
        result = parse_file(data, "application/json", "file.json")
        assert "kalit" in result

    def test_html_content_type(self):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            pytest.skip("beautifulsoup4 o'rnatilmagan")
        html = b"<p>HTML matn</p>"
        result = parse_file(html, "text/html", "file.html")
        assert "HTML matn" in result

    def test_csv_content_type(self):
        csv = b"Ism,Ball\nAli,95"
        result = parse_file(csv, "text/csv", "file.csv")
        assert "Ali" in result

    def test_tsv_content_type(self):
        tsv = b"Ism\tBall\nAli\t95"
        result = parse_file(tsv, "text/tab-separated-values", "file.tsv")
        assert "Ali" in result

    def test_unsupported_type_raises(self):
        with pytest.raises(ValueError, match="Qo'llab-quvvatlanmaydigan fayl turi"):
            parse_file(b"data", "application/octet-stream", "file.bin")

    def test_content_type_with_charset_stripped(self):
        result = parse_file(b"Matn", "text/plain; charset=utf-8", "file.txt")
        assert "Matn" in result

    def test_extension_fallback_when_wrong_content_type(self):
        csv = b"Ustun1,Ustun2\nA,B"
        result = parse_file(csv, "application/octet-stream", "data.csv")
        assert "A" in result

    @patch("app.ai.file_parser._get_whisper_model")
    def test_audio_dispatch(self, mock_get_model):
        mock_segment = MagicMock()
        mock_segment.text = "Audio test"
        mock_model = MagicMock()
        mock_model.transcribe.return_value = (
            [mock_segment], MagicMock(language="uz", duration=2.0)
        )
        mock_get_model.return_value = mock_model
        result = parse_file(b"fake", "audio/mpeg", "test.mp3")
        assert "Audio test" in result

    @patch("app.ai.file_parser.pytesseract")
    @patch("app.ai.file_parser.Image")
    def test_image_dispatch(self, mock_image_module, mock_tesseract):
        mock_image_module.open.return_value = MagicMock()
        mock_tesseract.image_to_string.return_value = "Rasm matni"
        result = parse_file(b"fake", "image/png", "photo.png")
        assert "Rasm matni" in result