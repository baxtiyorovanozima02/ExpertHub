"""
app/ai/file_parser.py

Fayl turlariga qarab matn chiqaradi:
  - PDF       → pdfplumber orqali
  - Word      → python-docx orqali
  - Excel     → openpyxl orqali
  - CSV/TSV   → csv module orqali
  - PowerPoint→ python-pptx orqali
  - TXT/MD/HTML/JSON/XML/RTF → to'g'ridan o'qiladi yoki parse qilinadi
  - EPUB      → ebooklib orqali
  - Rasm      → pytesseract (OCR) orqali
  - Audio     → faster-whisper (lokal, bepul) orqali
"""

import io
import json
import logging
import os
import tempfile
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber o'rnatilmagan: pip install pdfplumber")

    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())

    result = "\n\n".join(text_parts)
    if not result.strip():
        raise ValueError("PDF dan matn topilmadi (skanerlangan rasm bo'lishi mumkin)")
    return result



def extract_text_from_word(file_bytes: bytes) -> str:
    try:
        import docx as python_docx
    except ImportError:
        raise RuntimeError("python-docx o'rnatilmagan: pip install python-docx")

    doc = python_docx.Document(io.BytesIO(file_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("Word hujjatidan matn topilmadi")
    return text



def extract_text_from_excel(file_bytes: bytes, filename: Optional[str] = None) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl o'rnatilmagan: pip install openpyxl")

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    text_parts = []

    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        text_parts.append(f"## Sheet: {sheet.title}")

        header = [str(c).strip() if c is not None else "" for c in rows[0]]
        has_header = any(header)

        data_rows = rows[1:] if has_header else rows
        for row in data_rows:
            if row is None or all(c is None for c in row):
                continue
            if has_header:
                pairs = [
                    f"{header[i] or f'ustun{i+1}'}: {row[i]}"
                    for i in range(len(row))
                    if i < len(header) and row[i] is not None
                ]
            else:
                pairs = [f"ustun{i+1}: {v}" for i, v in enumerate(row) if v is not None]
            if pairs:
                text_parts.append(", ".join(pairs))

    result = "\n".join(text_parts).strip()
    if not result:
        raise ValueError("Excel fayldan ma'lumot topilmadi")
    return result


def extract_text_from_csv(file_bytes: bytes, delimiter: str = ",") -> str:
    import csv as csv_module

    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv_module.reader(io.StringIO(text), delimiter=delimiter)
    rows = list(reader)
    if not rows:
        raise ValueError("CSV fayl bo'sh")

    header = rows[0]
    text_parts = []
    for row in rows[1:]:
        pairs = [
            f"{header[i] if i < len(header) else f'ustun{i+1}'}: {val}"
            for i, val in enumerate(row)
            if val
        ]
        if pairs:
            text_parts.append(", ".join(pairs))

    result = "\n".join(text_parts).strip()
    if not result:
        raise ValueError("CSV fayldan ma'lumot topilmadi")
    return result


def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx o'rnatilmagan: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text_parts.append(f"## Slayd {i}\n" + "\n".join(slide_texts))

    result = "\n\n".join(text_parts).strip()
    if not result:
        raise ValueError("PowerPoint fayldan matn topilmadi")
    return result


def extract_text_from_epub(file_bytes: bytes) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError(
            "ebooklib yoki beautifulsoup4 o'rnatilmagan: "
            "pip install ebooklib beautifulsoup4"
        )

    with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        book = epub.read_epub(tmp_path)
        text_parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n").strip()
            if text:
                text_parts.append(text)
        result = "\n\n".join(text_parts).strip()
        if not result:
            raise ValueError("EPUB fayldan matn topilmadi")
        return result
    finally:
        os.unlink(tmp_path)



def extract_text_from_html(file_bytes: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("beautifulsoup4 o'rnatilmagan: pip install beautifulsoup4")

    soup = BeautifulSoup(file_bytes.decode("utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    if not text:
        raise ValueError("HTML fayldan matn topilmadi")
    return text


def extract_text_from_json(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace").strip()
    if "\n" in text:
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        parts = []
        for line in lines:
            try:
                obj = json.loads(line)
                parts.append(json.dumps(obj, ensure_ascii=False, indent=2))
            except Exception:
                parts.append(line)
        return "\n\n".join(parts)
    try:
        obj = json.loads(text)
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return text



def extract_text_from_xml(file_bytes: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes.decode("utf-8", errors="replace"), "xml")
        text = soup.get_text(separator="\n").strip()
        if text:
            return text
    except Exception:
        pass
    return file_bytes.decode("utf-8", errors="replace")



def extract_text_from_rtf(file_bytes: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(file_bytes.decode("utf-8", errors="replace"))
        if text.strip():
            return text.strip()
    except ImportError:
        pass

    import re
    raw = file_bytes.decode("utf-8", errors="replace")
    text = re.sub(r"\\[a-z]+\d*\s?", " ", raw)
    text = re.sub(r"[{}\\]", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        raise ValueError("RTF fayldan matn topilmadi")
    return text



def extract_text_from_image(file_bytes: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise RuntimeError(
            "pytesseract yoki Pillow o'rnatilmagan: pip install pytesseract Pillow"
        )

    image = Image.open(io.BytesIO(file_bytes))
    text = pytesseract.image_to_string(image, lang="uzb+rus+eng")
    if not text.strip():
        raise ValueError("Rasmdan matn topilmadi")
    return text.strip()



_whisper_model = None


def _get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        try:
            from faster_whisper import WhisperModel
        except ImportError:
            raise RuntimeError(
                "faster-whisper o'rnatilmagan: pip install faster-whisper"
            )
        logger.info("faster-whisper 'base' modeli yuklanmoqda...")
        _whisper_model = WhisperModel("base", device="cpu", compute_type="int8")
        logger.info("faster-whisper modeli tayyor.")
    return _whisper_model


def extract_text_from_audio(file_bytes: bytes, filename: str = "audio.mp3") -> str:
    ext = os.path.splitext(filename)[-1].lower() or ".mp3"

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        model = _get_whisper_model()
        segments, info = model.transcribe(tmp_path, beam_size=5)
        text = " ".join(segment.text for segment in segments).strip()
        if not text:
            raise ValueError("Audio dan matn topilmadi")
        logger.info(f"Audio til: {info.language}, davomiylik: {info.duration:.1f}s")
        return text
    finally:
        os.unlink(tmp_path)




def parse_file(
    file_bytes: bytes,
    content_type: str,
    filename: Optional[str] = None,
) -> str:
    ct = content_type.lower().split(";")[0].strip()  # charset ni olib tashlash

    if ct == "application/pdf":
        logger.info(f"PDF parsing: {filename}")
        return extract_text_from_pdf(file_bytes)

    if ct in ("text/plain", "text/markdown", "text/x-markdown", "text/log"):
        logger.info(f"Plain text: {filename}")
        return file_bytes.decode("utf-8", errors="replace")

    if ct in (
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        logger.info(f"Word parsing: {filename}")
        return extract_text_from_word(file_bytes)


    if ct in (
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ):
        logger.info(f"Excel parsing: {filename}")
        return extract_text_from_excel(file_bytes, filename)

    if ct in ("text/csv", "application/csv", "text/tab-separated-values"):
        delimiter = "\t" if ct == "text/tab-separated-values" else ","
        logger.info(f"CSV/TSV parsing: {filename}")
        return extract_text_from_csv(file_bytes, delimiter)

    if ct in (
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        logger.info(f"PPTX parsing: {filename}")
        return extract_text_from_pptx(file_bytes)

    # ── EPUB ─────────────────────────────────
    if ct == "application/epub+zip":
        logger.info(f"EPUB parsing: {filename}")
        return extract_text_from_epub(file_bytes)


    if ct in ("text/html", "application/xhtml+xml"):
        logger.info(f"HTML parsing: {filename}")
        return extract_text_from_html(file_bytes)


    if ct in ("application/json", "application/x-ndjson", "application/jsonlines",
              "text/json"):
        logger.info(f"JSON parsing: {filename}")
        return extract_text_from_json(file_bytes)

    if ct in ("application/xml", "text/xml"):
        logger.info(f"XML parsing: {filename}")
        return extract_text_from_xml(file_bytes)

    if ct == "application/rtf" or ct == "text/rtf":
        logger.info(f"RTF parsing: {filename}")
        return extract_text_from_rtf(file_bytes)


    if ct.startswith("image/"):
        logger.info(f"OCR (rasm): {filename}")
        return extract_text_from_image(file_bytes)


    if ct.startswith("audio/") or ct.startswith("video/"):
        logger.info(f"STT (audio): {filename}")
        return extract_text_from_audio(file_bytes, filename or "audio.mp3")


    if filename:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        ext_fallback = {
            "txt": lambda: file_bytes.decode("utf-8", errors="replace"),
            "md":  lambda: file_bytes.decode("utf-8", errors="replace"),
            "csv": lambda: extract_text_from_csv(file_bytes),
            "tsv": lambda: extract_text_from_csv(file_bytes, delimiter="\t"),
            "json": lambda: extract_text_from_json(file_bytes),
            "jsonl": lambda: extract_text_from_json(file_bytes),
            "xml": lambda: extract_text_from_xml(file_bytes),
            "html": lambda: extract_text_from_html(file_bytes),
            "htm":  lambda: extract_text_from_html(file_bytes),
            "rtf":  lambda: extract_text_from_rtf(file_bytes),
            "docx": lambda: extract_text_from_word(file_bytes),
            "doc":  lambda: extract_text_from_word(file_bytes),
            "xlsx": lambda: extract_text_from_excel(file_bytes, filename),
            "xls":  lambda: extract_text_from_excel(file_bytes, filename),
            "pptx": lambda: extract_text_from_pptx(file_bytes),
            "epub": lambda: extract_text_from_epub(file_bytes),
            "pdf":  lambda: extract_text_from_pdf(file_bytes),
        }
        if ext in ext_fallback:
            logger.info(f"Kengaytma bo'yicha fallback: .{ext} → {filename}")
            return ext_fallback[ext]()

    raise ValueError(f"Qo'llab-quvvatlanmaydigan fayl turi: {content_type}")